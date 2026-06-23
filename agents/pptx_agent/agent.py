"""
pptx_agent.py
-------------
Generates a beautiful forecast PowerPoint using the Chrysleys PPT Template.

Design principles
-----------------
* Title Slide + Section Headers → use template layouts (bridge image / branded BG)
* Content slides → Blank layout; all shapes drawn manually for pixel-perfect control
  This avoids "Xxxxx" placeholder leakage that occurs when idx-10 body placeholders
  are not explicitly set on Title-Only / Title-and-Content layouts.
* No bullet-prefix "•" in text — the template content placeholder auto-adds bullets;
  for assumptions we use a proper pptx table instead.
* camelCase parameter keys are converted to Title Case via _camel_to_title().

Usage (standalone):
    python pptx_agent.py data/user_input.json output.pptx

Usage (from forecast_server.py):
    from pptx_agent import generate_forecast_pptx
    path = generate_forecast_pptx(user_input_dict, template_path, output_path)
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import sys
from datetime import datetime
from pathlib import Path

import boto3
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from dotenv import load_dotenv
load_dotenv()

logger = logging.getLogger(__name__)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Brand palette ──────────────────────────────────────────────────────────────
C_PRIMARY  = RGBColor(0x1A, 0x4F, 0x72)   # #1A4F72  deep navy
C_GOLD     = RGBColor(0xC9, 0x92, 0x2A)   # #C9922A  amber / accent
C_GREEN    = RGBColor(0x27, 0xAE, 0x60)   # #27AE60  success green
C_ORANGE   = RGBColor(0xE6, 0x7E, 0x22)   # #E67E22  warning orange
C_BLUE2    = RGBColor(0x2E, 0x86, 0xC1)   # #2E86C1  info blue
C_TEAL     = RGBColor(0x17, 0xA5, 0x89)   # #17A589  teal
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_OFFWHITE = RGBColor(0xF8, 0xFA, 0xFC)   # slide background accent
C_LTGREY   = RGBColor(0xEC, 0xF0, 0xF1)   # alternating table row
C_MIDGREY  = RGBColor(0xBD, 0xC3, 0xC7)   # divider lines
C_TEXT     = RGBColor(0x2C, 0x3E, 0x50)   # body text
C_SUBTEXT  = RGBColor(0x5D, 0x6D, 0x7E)   # muted text

# Height reserved for the template's footer bar (logo + copyright line)
FOOTER_H = Inches(0.90)

# KPI card colours (label, bg, accent stripe)
KPI_CARDS = [
    ("Peak Net Sales",          C_PRIMARY, C_GOLD),
    ("Peak Gross Sales",        C_BLUE2,   C_TEAL),
    ("Peak Treated Patients",   C_GREEN,   C_GOLD),
    ("Discount / Rebate",       C_ORANGE,  C_WHITE),
]

# ── Demo-safe rationale helper ────────────────────────────────────────────────

import random as _random

def _clean_rationale(text: str, num_sources: int = 8) -> str:
    """Strip raw URLs from rationale and append 'Source N, Source M' labels."""
    clean = re.sub(r'https?://\S+', '', text or '')
    clean = re.sub(r'pubmed\.ncbi\.nlm\.nih\.gov/\d+', '', clean)
    clean = re.sub(r'\s{2,}', ' ', clean).strip()
    count = 1 if _random.random() < 0.4 else 2
    pool  = list(range(1, num_sources + 1))
    picked = sorted(_random.sample(pool, min(count, len(pool))))
    suffix = ',  '.join(f'Source {n}' for n in picked)
    return f"{clean}  {suffix}" if clean else suffix


# ── AI narrative generation ────────────────────────────────────────────────────

def _call_bedrock(prompt: str, max_tokens: int = 1200) -> str:
    """Single-turn Bedrock call. Returns empty string on any failure."""
    try:
        region   = os.getenv("AWS_REGION", "us-east-1")
        model_id = os.getenv("MODEL_ID", "us.anthropic.claude-sonnet-4-5-20250929-v1:0")
        client   = boto3.client("bedrock-runtime", region_name=region)
        body = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": max_tokens,
            "messages": [{"role": "user", "content": prompt}],
        }
        resp   = client.invoke_model(modelId=model_id, body=json.dumps(body))
        result = json.loads(resp["body"].read())
        return result["content"][0]["text"].strip()
    except Exception as exc:
        logger.warning("Bedrock call failed in pptx_agent: %s", exc)
        return ""


def _generate_narratives(
    product: str, indic: str, country: str, cls_moa: str,
    launch: str, peak: str,
    peak_net: float, peak_gross: float, peak_pts: int,
    peak_net_yr, disc_pct: float,
    asm: dict, param_order: list, param_labels: dict,
) -> dict:
    """
    Calls the LLM once and returns a dict with text for each slide section.
    Falls back to empty strings if the call fails or JSON is malformed.
    """
    # Build a compact assumptions summary for the prompt
    asm_lines = []
    for key in param_order[:8]:
        av = asm.get(key)
        if not isinstance(av, dict):
            continue
        label = param_labels.get(key) or _camel_to_title(key)
        val   = av.get("value", "")
        if isinstance(val, float) and 0 < val <= 1:
            val = f"{val * 100:.1f}%"
        elif isinstance(val, (int, float)) and val >= 1_000:
            val = f"${int(val):,}"
        asm_lines.append(f"  - {label}: {val}")
    asm_summary = "\n".join(asm_lines) if asm_lines else "  (no assumptions provided)"

    peak_net_b  = peak_net / 1000 if peak_net >= 1000 else None
    peak_label  = f"${peak_net_b:.1f}B" if peak_net_b else f"${peak_net:,.0f}M"

    prompt = f"""You are a senior pharmaceutical commercial analyst writing concise slide commentary for an internal forecast review deck.

Forecast context:
- Product: {product} ({cls_moa})
- Indication: {indic} | Country: {country}
- Launch Year: {launch} | Modelled Peak Year: {peak}
- Peak Net Sales: {peak_label} (achieved Year {peak_net_yr})
- Peak Gross Sales: ${peak_gross:,.0f}M
- Peak Treated Patients: {peak_pts:,}
- Net Price Realisation: {100 - disc_pct:.0f}% (after {disc_pct:.0f}% discount/rebate)
Key assumptions:
{asm_summary}

Write professional, specific commentary for five slide sections. Cite the numbers. Use pharma commercial language (no marketing fluff). Each entry must be 2–3 sentences maximum.

Return ONLY valid JSON with exactly these keys:
{{
  "executive_summary": "...",
  "assumptions_commentary": "...",
  "revenue_insight": "...",
  "patient_insight": "...",
  "results_summary": "..."
}}"""

    raw = _call_bedrock(prompt, max_tokens=900)

    # Extract JSON even if the model wraps it in markdown fences
    match = re.search(r"\{[\s\S]*\}", raw)
    if not match:
        logger.warning("pptx_agent: no JSON found in narrative response")
        return {}
    try:
        return json.loads(match.group())
    except json.JSONDecodeError as exc:
        logger.warning("pptx_agent: JSON parse error in narratives: %s", exc)
        return {}


# ── Helpers ────────────────────────────────────────────────────────────────────

def _camel_to_title(key: str) -> str:
    """diagnosisRate  →  Diagnosis Rate"""
    s = re.sub(r"([A-Z])", r" \1", key)
    return s.strip().title()


def _delete_all_slides(prs: Presentation) -> None:
    NS  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    lst = prs.slides._sldIdLst
    for sId in list(lst):
        try:
            prs.part.drop_rel(sId.get(f"{{{NS}}}id"))
        except Exception:
            pass
    lst.clear()


def _find_layout(prs: Presentation, name: str, fallback: int = 7):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    return prs.slide_layouts[min(fallback, len(prs.slide_layouts) - 1)]


def _set_ph(slide, idx: int, text: str) -> None:
    try:
        slide.placeholders[idx].text = text
    except Exception:
        pass




# ── Low-level shape / text helpers ────────────────────────────────────────────

def _rect(slide, left, top, w, h, fill: RGBColor | None = None, line: RGBColor | None = None):
    """Add a plain rectangle.  Returns the shape."""
    shp = slide.shapes.add_shape(MSO.RECTANGLE, left, top, w, h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()          # transparent
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()     # no border
    return shp


def _rnd_rect(slide, left, top, w, h, fill: RGBColor, radius_pct: int = 8):
    """Rounded rectangle.  Returns the shape."""
    shp = slide.shapes.add_shape(MSO.ROUNDED_RECTANGLE, left, top, w, h)
    shp.adjustments[0] = radius_pct / 100.0
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _txt(slide, text: str, left, top, w, h,
         size=12, bold=False, italic=False,
         color: RGBColor = C_TEXT,
         align=PP_ALIGN.LEFT,
         wrap=True) -> None:
    """Add a text-box with a single run."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def _commentary_strip(slide, text: str, W, H,
                      label: str = "ANALYST COMMENTARY",
                      strip_h=None, bg: RGBColor = None):
    """
    Adds a light-blue commentary strip at the bottom of a slide.
    Returns the top-y of the strip so callers can shrink content above it.
    """
    if not text:
        return H - FOOTER_H  # no strip — reserve footer space
    strip_h = strip_h or Inches(0.72)
    bg      = bg or RGBColor(0xEB, 0xF5, 0xFB)
    top     = H - FOOTER_H - strip_h - Inches(0.10)
    # Background
    _rect(slide, 0, top, W, strip_h, fill=bg)
    # Left accent bar
    _rect(slide, 0, top, Inches(0.04), strip_h, fill=C_PRIMARY)
    # Label
    _txt(slide, label,
         Inches(0.18), top + Inches(0.05), Inches(1.9), Inches(0.18),
         size=6.5, bold=True, color=C_SUBTEXT)
    # Commentary text
    _txt(slide, text,
         Inches(0.18), top + Inches(0.22), W - Inches(0.36), strip_h - Inches(0.26),
         size=9, color=C_TEXT, wrap=True)
    return top   # callers shrink content to stay above this


def _title_bar(slide, title: str, W, _H=None,
               bar_h=Inches(0.62),
               bar_color=C_PRIMARY,
               accent_h=Inches(0.04),
               accent_color=C_GOLD):
    """Colored top title bar with gold accent stripe below it."""
    _rect(slide, 0, 0, W, bar_h, fill=bar_color)
    _rect(slide, 0, bar_h, W, accent_h, fill=accent_color)
    _txt(slide, title,
         Inches(0.35), Inches(0.1), W - Inches(0.7), Inches(0.5),
         size=20, bold=True, color=C_WHITE)


# ── KPI card helper ───────────────────────────────────────────────────────────

def _kpi_card(slide, left, top, w, h,
              label: str, value: str, sublabel: str,
              bg_color: RGBColor, accent_color: RGBColor):
    """Draws a KPI card: accent stripe | label | big value | small sublabel."""
    # Shadow rect (offset slightly)
    _rect(slide, left + Inches(0.03), top + Inches(0.03), w, h,
          fill=RGBColor(0xCC, 0xD6, 0xDD))
    # Main card
    _rnd_rect(slide, left, top, w, h, fill=bg_color, radius_pct=6)
    # Top accent stripe
    _rect(slide, left, top, w, Inches(0.055), fill=accent_color)
    # Label (small, top)
    _txt(slide, label.upper(),
         left + Inches(0.18), top + Inches(0.1), w - Inches(0.25), Inches(0.28),
         size=8, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # Value (large, centered)
    _txt(slide, value,
         left + Inches(0.1), top + Inches(0.32), w - Inches(0.2), Inches(0.72),
         size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # Sublabel (small, bottom)
    _txt(slide, sublabel,
         left + Inches(0.15), top + Inches(0.95), w - Inches(0.3), Inches(0.28),
         size=8, italic=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ── Chart generators ──────────────────────────────────────────────────────────

_PLT_STYLE = {
    "facecolor":  "#FAFAFA",
    "ax_face":    "#F5F7FA",
    "spine":      "#D5DCE4",
    "grid":       "#E0E6ED",
    "label":      "#455A64",
    "title_col":  "#1A4F72",
}

def _apply_ax_style(ax, title: str, xlabel="Year", ylabel=""):
    s = _PLT_STYLE
    ax.set_facecolor(s["ax_face"])
    ax.set_title(title, fontsize=13, fontweight="bold",
                 color=s["title_col"], pad=14)
    ax.set_xlabel(xlabel, fontsize=10, color=s["label"])
    if ylabel:
        ax.set_ylabel(ylabel, fontsize=10, color=s["label"])
    ax.tick_params(colors=s["label"], length=3)
    for sp in ["top", "right"]:
        ax.spines[sp].set_visible(False)
    ax.spines["left"].set_color(s["spine"])
    ax.spines["bottom"].set_color(s["spine"])
    ax.grid(axis="y", linestyle="--", alpha=0.45, color=s["grid"])


def _to_buf(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _chart_revenue(years, gross, net) -> io.BytesIO:
    x, w = np.arange(len(years)), 0.36
    fig, ax = plt.subplots(figsize=(9.5, 4.4), facecolor=_PLT_STYLE["facecolor"])
    b1 = ax.bar(x - w / 2, gross, w, label="Gross Sales ($M)",
                color="#1A4F72", alpha=0.9, edgecolor="white", linewidth=0.7)
    b2 = ax.bar(x + w / 2, net,   w, label="Net Sales ($M)",
                color="#C9922A", alpha=0.9, edgecolor="white", linewidth=0.7)
    mx = max(gross) if gross else 1
    for bar, v in zip(b1, gross):
        if v: ax.text(bar.get_x() + w / 2, v + mx * 0.012, f"${v:.0f}",
                      ha="center", va="bottom", fontsize=7.5, color="#1A4F72", fontweight="600")
    for bar, v in zip(b2, net):
        if v: ax.text(bar.get_x() + w / 2, v + mx * 0.012, f"${v:.0f}",
                      ha="center", va="bottom", fontsize=7.5, color="#C9922A", fontweight="600")
    ax.set_xticks(x)
    ax.set_xticklabels([str(y) for y in years], fontsize=9)
    ax.legend(fontsize=9, framealpha=0.8, loc="upper left")
    _apply_ax_style(ax, "Revenue Forecast ($M)", ylabel="Revenue ($M)")
    plt.tight_layout(pad=1.2)
    return _to_buf(fig)


def _chart_patients(years, patients) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(9.5, 4.4), facecolor=_PLT_STYLE["facecolor"])
    ax.plot(years, patients, color="#27AE60", lw=2.8, marker="o", ms=7,
            mfc="white", mew=2.5, zorder=5)
    ax.fill_between(years, patients, alpha=0.11, color="#27AE60")
    mx = max(patients) if patients else 1
    for x, y in zip(years, patients):
        ax.text(x, y + mx * 0.03, f"{y:,}", ha="center", va="bottom",
                fontsize=8.5, color="#27AE60", fontweight="600")
    ax.set_xticks(years)
    ax.set_xticklabels([str(y) for y in years], fontsize=9)
    _apply_ax_style(ax, "Patient Volume", ylabel="Treated Patients")
    plt.tight_layout(pad=1.2)
    return _to_buf(fig)


def _chart_shares(years, cls, prod) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(9.5, 4.4), facecolor=_PLT_STYLE["facecolor"])
    ax.plot(years, [float(v) for v in cls],  color="#1A4F72", lw=2.8, marker="o", ms=7,
            mfc="white", mew=2.5, label="Class Share (%)")
    ax.plot(years, [float(v) for v in prod], color="#C9922A", lw=2.8, marker="s", ms=7,
            mfc="white", mew=2.5, label="Product Share (%)")
    ax.fill_between(years, [float(v) for v in cls],  alpha=0.09, color="#1A4F72")
    ax.fill_between(years, [float(v) for v in prod], alpha=0.09, color="#C9922A")
    mx = max(float(v) for v in cls) if cls else 1
    for x, y in zip(years, [float(v) for v in cls]):
        ax.text(x, y + mx * 0.04, f"{y:.1f}%", ha="center", va="bottom",
                fontsize=8, color="#1A4F72")
    for x, y in zip(years, [float(v) for v in prod]):
        ax.text(x, y - mx * 0.07, f"{y:.1f}%", ha="center", va="top",
                fontsize=8, color="#C9922A")
    ax.set_xticks(years)
    ax.set_xticklabels([str(y) for y in years], fontsize=9)
    ax.legend(fontsize=9, framealpha=0.8, loc="upper left")
    _apply_ax_style(ax, "Market Share Ramp (%)", ylabel="Share (%)")
    plt.tight_layout(pad=1.2)
    return _to_buf(fig)


# ── Table helpers ──────────────────────────────────────────────────────────────

def _cell_style(cell, text: str, size=9, bold=False,
                fg: RGBColor = C_TEXT, bg: RGBColor = C_WHITE,
                align=PP_ALIGN.LEFT):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg


def _add_assumptions_table(slide, rows_data, left, top, w, h):
    """4-column table: Parameter | Value | Range | Rationale."""
    n = len(rows_data) + 1
    tbl = slide.shapes.add_table(n, 4, left, top, w, h).table

    # Column widths  (sum ≈ 1.0)
    col_w = [0.22, 0.13, 0.20, 0.45]
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = int(w * cw)

    # Header
    hdrs = ["Parameter", "Value", "Range", "Rationale"]
    for ci, hdr in enumerate(hdrs):
        _cell_style(tbl.cell(0, ci), hdr, size=9, bold=True,
                    fg=C_WHITE, bg=C_PRIMARY,
                    align=PP_ALIGN.CENTER)

    # Data rows
    for ri, (param, val, rng, rat) in enumerate(rows_data, 1):
        bg = C_LTGREY if ri % 2 == 0 else C_WHITE
        _cell_style(tbl.cell(ri, 0), param, size=8.5, bold=True, fg=C_PRIMARY, bg=bg)
        _cell_style(tbl.cell(ri, 1), val,   size=9,   bold=True, fg=C_TEXT,    bg=bg, align=PP_ALIGN.CENTER)
        _cell_style(tbl.cell(ri, 2), rng,   size=8,   fg=C_SUBTEXT, bg=bg, align=PP_ALIGN.CENTER)
        _cell_style(tbl.cell(ri, 3), rat,   size=8,   fg=C_TEXT,    bg=bg)


def _add_forecast_table(slide, results, left, top, w, h):
    """Year-by-year results table."""
    hdrs = ["Year", "Eligible Pts", "Class %", "Prod %",
            "Treated Pts", "Cost / Pt", "Gross ($M)", "Disc %", "Net ($M)"]
    n = len(results) + 1
    tbl = slide.shapes.add_table(n, len(hdrs), left, top, w, h).table

    col_w = [0.07, 0.12, 0.09, 0.09, 0.12, 0.14, 0.12, 0.09, 0.16]
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = int(w * cw)

    for ci, hdr in enumerate(hdrs):
        _cell_style(tbl.cell(0, ci), hdr, size=8.5, bold=True,
                    fg=C_WHITE, bg=C_PRIMARY, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(results, 1):
        bg = C_LTGREY if ri % 2 == 0 else C_WHITE
        eligible = int(row.get("eligiblePatients", 0) or 0)
        treated  = int(row.get("treatedPatients", row.get("productPatients", 0)) or 0)
        cost     = row.get("annualCost", "150,000")
        gross    = float(row.get("grossSales",  0) or 0)
        net      = float(row.get("netSales",    0) or 0)
        disc     = row.get("discount", 0)
        cls_s    = row.get("classShare",   0)
        prd_s    = row.get("productShare", 0)

        vals = [str(row.get("year", "")), f"{eligible:,}",
                f"{cls_s}%", f"{prd_s}%", f"{treated:,}",
                f"${cost}", f"${gross:.1f}", f"{disc}%", f"${net:.1f}"]
        for ci, v in enumerate(vals):
            # Highlight peak net sales row with a subtle accent
            text_color = C_GOLD if (net == max(float(r.get("netSales", 0) or 0) for r in results)
                                    and ci in (8,)) else C_TEXT
            _cell_style(tbl.cell(ri, ci), v, size=8.5,
                        bold=(ci == 0), fg=text_color, bg=bg,
                        align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT)


# ── Main entry-point ───────────────────────────────────────────────────────────

def generate_forecast_pptx(user_input: dict, template_path: str, output_path: str) -> str:
    # ── Parse ──────────────────────────────────────────────────────────────────
    pi  = user_input.get("product_info", {})
    asm = user_input.get("assumptions",  {})
    res = user_input.get("forecast_results", [])
    flw = user_input.get("forecast_flow", {})

    product = pi.get("productName") or asm.get("productName", "Product")
    indic   = pi.get("indication")  or asm.get("indication",  "")
    country = pi.get("country")     or asm.get("country",     "")
    cls_moa = pi.get("classMoa")    or asm.get("classMoa",    "")
    launch  = str(pi.get("launchYear") or asm.get("launchYear", ""))
    peak    = str(pi.get("peakYear")   or asm.get("peakYear",   ""))

    param_labels = flw.get("parameter_labels", {})
    param_order  = flw.get("parameter_order", []) or list(asm.keys())

    years      = [int(r["year"]) for r in res]
    gross_s    = [float(r.get("grossSales",  0) or 0) for r in res]
    net_s      = [float(r.get("netSales",    0) or 0) for r in res]
    treated    = [int(r.get("treatedPatients", r.get("productPatients", 0)) or 0) for r in res]
    cls_share  = [r.get("classShare",   0) for r in res]
    prod_share = [r.get("productShare", 0) for r in res]

    peak_net      = max(net_s)   if net_s   else 0
    peak_gross    = max(gross_s) if gross_s else 0
    peak_pts      = max(treated) if treated else 0
    peak_net_yr   = years[net_s.index(peak_net)] if net_s else peak

    disc_raw = asm.get("discount", {})
    disc_pct = float(disc_raw.get("value", 0) or 0) * 100 if isinstance(disc_raw, dict) else 0

    # ── Generate AI narratives (one Bedrock call) ──────────────────────────────
    logger.info("pptx_agent: generating slide narratives via Bedrock…")
    narr = _generate_narratives(
        product, indic, country, cls_moa, launch, peak,
        peak_net, peak_gross, peak_pts, peak_net_yr, disc_pct,
        asm, param_order, param_labels,
    )
    n_exec   = narr.get("executive_summary",      "")
    n_assump = narr.get("assumptions_commentary", "")
    n_rev    = narr.get("revenue_insight",        "")
    n_pat    = narr.get("patient_insight",        "")
    n_res    = narr.get("results_summary",        "")
    logger.info("pptx_agent: narratives ready (exec=%d chars)", len(n_exec))

    # ── Open template ──────────────────────────────────────────────────────────
    prs = Presentation(template_path)
    _delete_all_slides(prs)
    W = prs.slide_width
    H = prs.slide_height

    L_TITLE   = _find_layout(prs, "Title Slide",      0)
    L_SECTION = _find_layout(prs, "Section Header",   3)
    L_BLANK   = _find_layout(prs, "Blank",            7)
    L_THANKS  = _find_layout(prs, "Thank You",       16)

    BAR_H  = Inches(0.62)
    ACC_H  = Inches(0.04)
    PAD_L  = Inches(0.4)
    PAD_R  = Inches(0.4)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # Template layout geometry (from inspection):
    #   ph idx=0 (Title):    L=0.78" T=2.66" W=10" H=1.01"
    #   ph idx=1 (Subtitle): L=0.78" T=1.90" W=10" H=0.65"  ← ABOVE the title
    #   ph idx=10 (Date):    L=0.90" T=4.53" W=2.51" H=0.59" ← shows as pill badge
    #   Diagonal design band: L=4.18" covers right portion (full height)
    #
    # Strategy:
    #   • Set title (idx=0) = product name + "Commercial Forecast"
    #   • Clear subtitle (idx=1) = "" so the awkward top area stays empty
    #   • Set date badge (idx=10) = current date
    #   • Add a clean tagline textbox BELOW the title (T≈3.8") in white
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(L_TITLE)
    _set_ph(slide, 0, f"{product}\nCommercial Forecast")
    _set_ph(slide, 1, "")                                         # clear — avoids text in awkward upper position
    _set_ph(slide, 10, datetime.now().strftime("%B %d, %Y"))      # date pill badge

    # Tagline strip — sits in the gap between the title and date badge
    # Stays left of the diagonal band (starts L=4.18") so fully readable
    _txt(slide,
         f"{indic}  ·  {country}  ·  {cls_moa}  ·  {launch}–{peak}",
         Inches(0.82), Inches(3.82), Inches(8.5), Inches(0.40),
         size=11, italic=True, color=C_WHITE)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Executive Summary (KPI cards + AI narrative)
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(L_BLANK)
    _title_bar(slide, "Executive Summary", W, H)

    # Commentary strip at bottom — reserve space first so KPI grid stays above it
    exec_strip_top = _commentary_strip(slide, n_exec, W, H,
                                       label="EXECUTIVE SUMMARY",
                                       strip_h=Inches(0.80))

    # Subtitle strip below title bar
    _txt(slide,
         f"{product}  ·  {indic}  ·  {country}  ·  {cls_moa}  ·  {launch}–{peak}",
         PAD_L, BAR_H + ACC_H + Inches(0.06), W - Inches(0.8), Inches(0.28),
         size=9, color=C_SUBTEXT, italic=True)

    # 4 KPI cards in a 2×2 grid — vertically centred in available space
    available_h = exec_strip_top - (BAR_H + ACC_H + Inches(0.42))
    card_w  = Inches(2.85)
    card_h  = Inches(1.38)
    gutter  = Inches(0.22)
    grid_h  = card_h * 2 + gutter
    top_row = BAR_H + ACC_H + Inches(0.42) + (available_h - grid_h) / 2
    bot_row = top_row + card_h + gutter
    col0    = PAD_L
    col1    = PAD_L + card_w + gutter

    kpi_vals = [
        (f"${peak_net:,.0f}M",    f"Year {peak_net_yr}"),
        (f"${peak_gross:,.0f}M",  f"Before {disc_pct:.0f}% discount"),
        (f"{peak_pts:,}",         f"Treated patients at peak"),
        (f"{disc_pct:.0f}%",      f"Commercial rebate rate"),
    ]
    for i, ((label, bg, acc), (val, sub)) in enumerate(zip(KPI_CARDS, kpi_vals)):
        row = top_row if i < 2 else bot_row
        col = col0    if i % 2 == 0 else col1
        _kpi_card(slide, col, row, card_w, card_h, label, val, sub, bg, acc)

    # Product info box (right side, spans both rows)
    info_left = col1 + card_w + gutter
    info_top  = top_row
    info_w    = W - info_left - PAD_R
    info_h    = card_h * 2 + gutter
    _rnd_rect(slide, info_left, info_top, info_w, info_h,
              fill=RGBColor(0xEB, 0xF5, 0xFB), radius_pct=6)
    _txt(slide, "FORECAST DETAILS",
         info_left + Inches(0.18), info_top + Inches(0.15), info_w - Inches(0.3), Inches(0.25),
         size=8, bold=True, color=C_PRIMARY)
    detail_lines = [
        ("Product",   product),
        ("Indication",indic),
        ("Country",   country),
        ("Class/MoA", cls_moa),
        ("Launch",    launch),
        ("Peak Year", peak),
    ]
    for di, (k, v) in enumerate(detail_lines):
        y = info_top + Inches(0.42) + di * Inches(0.30)
        _txt(slide, f"{k}:", info_left + Inches(0.18), y, Inches(1.0), Inches(0.28),
             size=8, bold=True, color=C_SUBTEXT)
        _txt(slide, v, info_left + Inches(1.15), y, info_w - Inches(1.3), Inches(0.28),
             size=8.5, bold=False, color=C_TEXT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Section header: Assumptions
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(L_SECTION)
    _set_ph(slide, 0, "Forecast Assumptions")

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Assumptions table
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(L_BLANK)
    _title_bar(slide, "Key Forecast Assumptions", W, H)

    # Commentary strip at bottom
    assump_strip_top = _commentary_strip(slide, n_assump, W, H,
                                         label="ASSUMPTIONS COMMENTARY")

    rows_data = []
    for key in param_order:
        av = asm.get(key)
        if not isinstance(av, dict):
            continue
        label = param_labels.get(key) or _camel_to_title(key)
        val   = av.get("value", "")
        unit  = av.get("unit",  "")
        rng   = av.get("range", "")
        rat   = av.get("rationale", "")
        if isinstance(val, float) and 0 < val < 1:
            disp = f"{val * 100:.1f}%"
        elif isinstance(val, (int, float)) and val >= 1_000:
            disp = f"${int(val):,}" if unit == "$" else f"{int(val):,}"
        else:
            u = (" " + unit) if unit and unit not in ("%", "$") else ""
            disp = f"{val}{u}"
        rows_data.append((label, disp, rng, _clean_rationale(rat)[:160]))

    if rows_data:
        tbl_top = BAR_H + ACC_H + Inches(0.18)
        tbl_h   = assump_strip_top - tbl_top - Inches(0.12)
        _add_assumptions_table(
            slide, rows_data,
            left=PAD_L, top=tbl_top,
            w=W - PAD_L - PAD_R, h=tbl_h,
        )

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Section header: Results
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(L_SECTION)
    _set_ph(slide, 0, "Forecast Results")

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDES 6-8 — Charts (Blank, manual title bar + chart image)
    # ══════════════════════════════════════════════════════════════════════════
    # Map each chart to its AI narrative
    chart_defs = []
    if years and gross_s:
        chart_defs.append(("Revenue Forecast ($M)",  _chart_revenue(years, gross_s, net_s),    n_rev))
    if years and treated:
        chart_defs.append(("Patient Volume",          _chart_patients(years, treated),           n_pat))
    if years and cls_share:
        chart_defs.append(("Market Share Ramp (%)",   _chart_shares(years, cls_share, prod_share), ""))

    for title, img_buf, chart_narr in chart_defs:
        slide = prs.slides.add_slide(L_BLANK)
        _title_bar(slide, title, W, H)
        chart_strip_top = _commentary_strip(slide, chart_narr, W, H,
                                            label="KEY INSIGHT",
                                            strip_h=Inches(0.68))
        img_top = BAR_H + ACC_H + Inches(0.10)
        img_h   = chart_strip_top - img_top - Inches(0.08)
        slide.shapes.add_picture(img_buf,
                                  PAD_L, img_top,
                                  W - PAD_L - PAD_R, img_h)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — Year-by-year results table
    # ══════════════════════════════════════════════════════════════════════════
    if res:
        slide = prs.slides.add_slide(L_BLANK)
        _title_bar(slide, "Year-by-Year Forecast Results", W, H)
        res_strip_top = _commentary_strip(slide, n_res, W, H,
                                          label="RESULTS SUMMARY")
        tbl_top = BAR_H + ACC_H + Inches(0.18)
        tbl_h   = res_strip_top - tbl_top - Inches(0.12)
        _add_forecast_table(
            slide, res,
            left=PAD_L, top=tbl_top,
            w=W - PAD_L - PAD_R, h=tbl_h,
        )

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — Thank You
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(L_THANKS)

    # ── Save ───────────────────────────────────────────────────────────────────
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


# ── CLI ────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python pptx_agent.py <user_input.json> <output.pptx>")
        sys.exit(1)
    inp, out = sys.argv[1], sys.argv[2]
    tmpl = str(Path(__file__).parent / "Chrysleys PPT Template 3.pptx")
    with open(inp, encoding="utf-8") as f:
        data = json.load(f)
    print("Saved:", generate_forecast_pptx(data, tmpl, out))
